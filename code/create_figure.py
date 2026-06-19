"""
Generate Figure 1: 5-Fold Cross-Validation for Residual-on-Residual Regression

This script creates an editable PowerPoint figure showing:
1. Dataset split into 5 folds
2. Cross-validation to get out-of-sample predictions for m_Y(C) and m_A(C)
3. Residual construction
4. Final residual-on-residual regression

Requirements: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as RgbColor

# Create presentation with wide layout
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Add blank slide
blank_layout = prs.slide_layouts[6]  # Blank layout
slide = prs.slides.add_slide(blank_layout)

# Color scheme
FOLD_TRAIN_COLOR = RgbColor(200, 220, 240)  # Light blue for training
FOLD_TEST_COLOR = RgbColor(255, 200, 100)   # Orange for test/validation
ARROW_COLOR = RgbColor(80, 80, 80)          # Dark gray arrows
MODEL_Y_COLOR = RgbColor(180, 220, 180)     # Light green for Y model
MODEL_A_COLOR = RgbColor(220, 180, 220)     # Light purple for A model
RESIDUAL_COLOR = RgbColor(255, 230, 180)    # Light yellow for residuals
FINAL_COLOR = RgbColor(180, 200, 230)       # Light blue for final regression

def add_text_box(slide, left, top, width, height, text, font_size=10, bold=False,
                 fill_color=None, align='center', font_name='Arial'):
    """Add a text box with optional fill color."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.text = text

    # Format text
    tf = shape.text_frame
    tf.word_wrap = True
    for para in tf.paragraphs:
        para.font.size = Pt(font_size)
        para.font.bold = bold
        para.font.name = font_name
        if align == 'center':
            para.alignment = PP_ALIGN.CENTER
        elif align == 'left':
            para.alignment = PP_ALIGN.LEFT

    # Vertical centering
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.paragraphs[0].font.size = Pt(font_size)

    # Fill color
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()

    # Border
    shape.line.color.rgb = RgbColor(100, 100, 100)
    shape.line.width = Pt(1)

    return shape

def add_fold_row(slide, left, top, fold_num, test_fold_idx):
    """Add a row showing one CV iteration with 5 folds."""
    fold_width = Inches(0.5)
    fold_height = Inches(0.35)
    gap = Inches(0.05)

    # Add fold number label
    label = slide.shapes.add_textbox(left - Inches(0.4), top, Inches(0.35), fold_height)
    tf = label.text_frame
    tf.text = f"{fold_num}"
    tf.paragraphs[0].font.size = Pt(9)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Add 5 folds
    for i in range(5):
        x = left + i * (fold_width + gap)
        if i == test_fold_idx:
            color = FOLD_TEST_COLOR
            text = "Test"
        else:
            color = FOLD_TRAIN_COLOR
            text = "Train"

        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, top, fold_width, fold_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RgbColor(80, 80, 80)
        shape.line.width = Pt(0.75)

        # Add text
        tf = shape.text_frame
        tf.text = text
        tf.paragraphs[0].font.size = Pt(7)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    return left + 5 * (fold_width + gap)

def add_arrow(slide, start_left, start_top, end_left, end_top, text=""):
    """Add a connector arrow."""
    # Use a line connector
    connector = slide.shapes.add_connector(
        1,  # straight connector
        start_left, start_top,
        end_left, end_top
    )
    connector.line.color.rgb = ARROW_COLOR
    connector.line.width = Pt(1.5)

def add_simple_arrow(slide, left, top, width, direction='right'):
    """Add a simple arrow shape."""
    if direction == 'right':
        shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, Inches(0.25))
    elif direction == 'down':
        shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, Inches(0.25), width)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ARROW_COLOR
    shape.line.fill.background()
    return shape

# =============================================================================
# LAYOUT
# =============================================================================

# Title
title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.5))
tf = title.text_frame
tf.text = "Figure 1. 5-Fold Cross-Validation for Residual-on-Residual Regression"
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.bold = True

# -----------------------------------------------------------------------------
# SECTION A: Dataset and 5-fold CV diagram
# -----------------------------------------------------------------------------

# Section label
section_a = slide.shapes.add_textbox(Inches(0.3), Inches(0.7), Inches(1), Inches(0.3))
tf = section_a.text_frame
tf.text = "A"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.bold = True

# "Dataset" box
dataset_box = add_text_box(slide, Inches(0.5), Inches(1.0), Inches(1.0), Inches(2.0),
                           "Dataset\n(n observations)\n\nY, A, C", font_size=10, bold=False,
                           fill_color=RgbColor(230, 230, 230))

# Arrow from dataset to folds
add_simple_arrow(slide, Inches(1.55), Inches(1.9), Inches(0.4), 'right')

# 5-fold CV iterations label
cv_label = slide.shapes.add_textbox(Inches(2.0), Inches(0.85), Inches(3.5), Inches(0.25))
tf = cv_label.text_frame
tf.text = "5-Fold Cross-Validation (k = 1, ..., 5)"
tf.paragraphs[0].font.size = Pt(10)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Iteration label
iter_label = slide.shapes.add_textbox(Inches(2.0), Inches(1.05), Inches(0.4), Inches(0.25))
tf = iter_label.text_frame
tf.text = "Iter."
tf.paragraphs[0].font.size = Pt(8)
tf.paragraphs[0].font.bold = True

# Fold labels (1-5)
for i in range(5):
    fold_label = slide.shapes.add_textbox(Inches(2.45) + i * Inches(0.55), Inches(1.05),
                                          Inches(0.5), Inches(0.2))
    tf = fold_label.text_frame
    tf.text = f"Fold {i+1}"
    tf.paragraphs[0].font.size = Pt(7)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Add 5 CV iteration rows
cv_start_top = Inches(1.3)
for k in range(5):
    add_fold_row(slide, Inches(2.45), cv_start_top + k * Inches(0.42), k+1, k)

# -----------------------------------------------------------------------------
# SECTION B: Fit models and get predictions
# -----------------------------------------------------------------------------

# Arrow from CV to models
add_simple_arrow(slide, Inches(5.3), Inches(1.9), Inches(0.5), 'right')

# "Fit on Train" label
fit_label = slide.shapes.add_textbox(Inches(5.35), Inches(1.5), Inches(0.5), Inches(0.4))
tf = fit_label.text_frame
tf.text = "Fit on\nTrain"
tf.paragraphs[0].font.size = Pt(7)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Section label B
section_b = slide.shapes.add_textbox(Inches(5.8), Inches(0.7), Inches(1), Inches(0.3))
tf = section_b.text_frame
tf.text = "B"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.bold = True

# Model m_Y(C) box
model_y_box = add_text_box(slide, Inches(5.9), Inches(1.1), Inches(1.3), Inches(0.7),
                           "Fit model\nm̂Y(C)\n(Y ~ C)", font_size=9,
                           fill_color=MODEL_Y_COLOR)

# Model m_A(C) box
model_a_box = add_text_box(slide, Inches(5.9), Inches(2.0), Inches(1.3), Inches(0.7),
                           "Fit model\nm̂A(C)\n(A ~ C)", font_size=9,
                           fill_color=MODEL_A_COLOR)

# Arrow from models to predictions
add_simple_arrow(slide, Inches(7.25), Inches(1.9), Inches(0.4), 'right')

# "Predict on Test" label
pred_label = slide.shapes.add_textbox(Inches(7.25), Inches(1.5), Inches(0.5), Inches(0.4))
tf = pred_label.text_frame
tf.text = "Predict\non Test"
tf.paragraphs[0].font.size = Pt(7)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Predictions boxes
pred_y_box = add_text_box(slide, Inches(7.7), Inches(1.1), Inches(1.3), Inches(0.7),
                          "Out-of-sample\npredictions\nŶ = m̂Y(C)", font_size=9,
                          fill_color=MODEL_Y_COLOR)

pred_a_box = add_text_box(slide, Inches(7.7), Inches(2.0), Inches(1.3), Inches(0.7),
                          "Out-of-sample\npredictions\nÂ = m̂A(C)", font_size=9,
                          fill_color=MODEL_A_COLOR)

# -----------------------------------------------------------------------------
# SECTION C: Construct residuals
# -----------------------------------------------------------------------------

# Arrow down to residuals
add_simple_arrow(slide, Inches(8.3), Inches(2.75), Inches(0.4), 'down')

# Section label C
section_c = slide.shapes.add_textbox(Inches(5.8), Inches(3.2), Inches(1), Inches(0.3))
tf = section_c.text_frame
tf.text = "C"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.bold = True

# "After all 5 iterations" label
combine_label = slide.shapes.add_textbox(Inches(5.9), Inches(3.45), Inches(3.2), Inches(0.35))
tf = combine_label.text_frame
tf.text = "Combine predictions from all 5 folds to get\nout-of-sample predictions for entire dataset"
tf.paragraphs[0].font.size = Pt(9)
tf.paragraphs[0].font.bold = False

# Residual Y box
resid_y_box = add_text_box(slide, Inches(5.9), Inches(3.95), Inches(1.5), Inches(0.65),
                           "Outcome Residual\nỸ = Y − Ŷ", font_size=10,
                           fill_color=RESIDUAL_COLOR)

# Residual A box
resid_a_box = add_text_box(slide, Inches(7.5), Inches(3.95), Inches(1.5), Inches(0.65),
                           "Exposure Residual\nÃ = A − Â", font_size=10,
                           fill_color=RESIDUAL_COLOR)

# -----------------------------------------------------------------------------
# SECTION D: Final regression
# -----------------------------------------------------------------------------

# Arrow down to final regression
add_simple_arrow(slide, Inches(7.25), Inches(4.65), Inches(0.5), 'down')

# Section label D
section_d = slide.shapes.add_textbox(Inches(5.8), Inches(5.15), Inches(1), Inches(0.3))
tf = section_d.text_frame
tf.text = "D"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.bold = True

# Final regression box
final_box = add_text_box(slide, Inches(5.9), Inches(5.4), Inches(3.1), Inches(0.8),
                         "Residual-on-Residual Regression\nỸ = ψ·Ã + ε\n(OLS with no intercept)",
                         font_size=11, bold=False, fill_color=FINAL_COLOR)

# Result box
result_box = add_text_box(slide, Inches(5.9), Inches(6.3), Inches(3.1), Inches(0.55),
                          "ψ̂ = confounder-adjusted effect estimate",
                          font_size=10, bold=True, fill_color=RgbColor(255, 255, 200))

# -----------------------------------------------------------------------------
# LEGEND
# -----------------------------------------------------------------------------

legend_top = Inches(4.5)
legend_left = Inches(0.5)

legend_title = slide.shapes.add_textbox(legend_left, legend_top, Inches(2), Inches(0.3))
tf = legend_title.text_frame
tf.text = "Legend"
tf.paragraphs[0].font.size = Pt(10)
tf.paragraphs[0].font.bold = True

# Legend items
legend_items = [
    (FOLD_TRAIN_COLOR, "Training folds"),
    (FOLD_TEST_COLOR, "Test/validation fold"),
    (MODEL_Y_COLOR, "Outcome model m̂Y(C)"),
    (MODEL_A_COLOR, "Exposure model m̂A(C)"),
    (RESIDUAL_COLOR, "Residuals"),
]

for i, (color, label) in enumerate(legend_items):
    y = legend_top + Inches(0.35) + i * Inches(0.3)

    # Color box
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, legend_left, y, Inches(0.3), Inches(0.2))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.color.rgb = RgbColor(100, 100, 100)
    box.line.width = Pt(0.75)

    # Label
    lbl = slide.shapes.add_textbox(legend_left + Inches(0.4), y, Inches(2), Inches(0.25))
    tf = lbl.text_frame
    tf.text = label
    tf.paragraphs[0].font.size = Pt(9)

# -----------------------------------------------------------------------------
# ANNOTATIONS / NOTES
# -----------------------------------------------------------------------------

notes_top = Inches(6.0)
notes_left = Inches(0.5)

notes = slide.shapes.add_textbox(notes_left, notes_top, Inches(4.5), Inches(1.2))
tf = notes.text_frame
tf.word_wrap = True
tf.text = ("Notes:\n"
           "• m̂Y(C) and m̂A(C) can use any flexible ML method (e.g., super learner)\n"
           "• Cross-validation ensures out-of-sample predictions avoid overfitting\n"
           "• Final OLS has no intercept because residuals are mean-zero by construction")
for para in tf.paragraphs:
    para.font.size = Pt(8)
tf.paragraphs[0].font.bold = True

# =============================================================================
# Save
# =============================================================================

output_path = "/Users/ain/Dropbox/01 Projects/AJE_Classroom/Residual_on_Residual/Figure1_RoR_CrossValidation.pptx"
prs.save(output_path)
print(f"Figure saved to: {output_path}")
print("\nYou can now open this in PowerPoint and modify colors, text, positions, etc.")
